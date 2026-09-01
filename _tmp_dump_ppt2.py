from pptx import Presentation

path = r"网易互娱通用PPT模板（含保密）_ExcelAgent答辩优化版_45分钟技术答辩充实版.pptx"
p = Presentation(path)
print("slides:", len(p.slides))
for i, s in enumerate(p.slides, 1):
    texts = []
    for shp in s.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip():
            texts.append(shp.text_frame.text.strip().splitlines()[0][:50])
    print(i, "|", " / ".join(texts[:2]))
