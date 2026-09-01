# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Emu

path = r"网易互娱通用PPT模板（含保密）_ExcelAgent答辩优化版.pptx"
p = Presentation(path)

def shape_sort_key(shp):
    try:
        return (round(shp.top / 91440), round(shp.left / 91440))
    except Exception:
        return (0, 0)

out = []
out.append(f"# PPTX 全文导出：{path}")
out.append(f"总页数：{len(p.slides)}")
out.append("")

for i, s in enumerate(p.slides, 1):
    out.append(f"\n{'='*80}\n## 第 {i} 页\n{'='*80}")
    shapes = list(s.shapes)
    try:
        shapes.sort(key=shape_sort_key)
    except Exception:
        pass
    for shp in shapes:
        if shp.has_table:
            tbl = shp.table
            out.append("[表格]")
            for row in tbl.rows:
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                out.append(" | ".join(cells))
            out.append("")
        elif shp.has_text_frame and shp.text_frame.text.strip():
            out.append(shp.text_frame.text.strip())
            out.append("")
        elif shp.shape_type == 6:  # group
            pass
    # notes
    if s.has_notes_slide:
        note_text = s.notes_slide.notes_text_frame.text.strip()
        if note_text:
            out.append(f"[备注] {note_text}")

with open("_current_ppt_dump.md", "w", encoding="utf-8") as f:
    f.write("\n".join(out))

print("done, lines:", len(out))
