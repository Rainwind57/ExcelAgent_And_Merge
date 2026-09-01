# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Pt

path = r"网易互娱通用PPT模板（含保密）_ExcelAgent答辩优化版.pptx"
p = Presentation(path)

def font_info(run):
    f = run.font
    color = None
    try:
        if f.color and f.color.type is not None:
            color = f.color.rgb
    except Exception:
        color = None
    return {
        "text": run.text[:20],
        "name": f.name,
        "size": f.size.pt if f.size else None,
        "bold": f.bold,
        "color": str(color) if color else None,
    }

sample_slides = [1, 2, 3, 4, 6, 7, 9, 20, 22, 24, 27, 29, 35, 36, 43]

for idx in sample_slides:
    s = p.slides[idx-1]
    print(f"\n===== 第{idx}页 =====")
    for shp in s.shapes:
        if shp.has_text_frame:
            tf = shp.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    info = font_info(run)
                    if info["text"].strip():
                        print(info)
        if shp.has_table:
            tbl = shp.table
            # header row + first data row
            for ridx, row in enumerate(tbl.rows):
                if ridx > 1:
                    break
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            info = font_info(run)
                            if info["text"].strip():
                                print("[table]", info)
